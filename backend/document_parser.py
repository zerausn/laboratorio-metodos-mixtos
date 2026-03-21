import os
import pandas as pd
from io import BytesIO

# Importar librerías específicas (envueltas en try-except por seguridad)
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import whisper
except ImportError:
    whisper = None

class DocumentParser:
    """Clase utilitaria para extraer texto de múltiples formatos de archivo."""
    
    @staticmethod
    def extract_text(file_obj, filename=""):
        """
        Recibe un objeto de archivo (ej. de st.file_uploader) o ruta de archivo,
        y extrae todo su texto basándose en la extensión.
        """
        ext = os.path.splitext(filename)[1].lower() if filename else ""
        
        # Leemos el contenido en Bytes si es necesario
        # Si file_obj es de Streamlit, tiene método read() y getvalue()
        
        texto_extraido = ""
        
        try:
            if ext == '.txt':
                texto_extraido = file_obj.getvalue().decode("utf-8")
                
            elif ext == '.pdf':
                if not fitz:
                    raise ImportError("Librería PyMuPDF (fitz) no instalada.")
                
                # Leemos el PDF desde el flujo de bytes de Streamlit
                pdf_bytes = file_obj.getvalue()
                with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                    paginas = []
                    for page in doc:
                        # get_text("text") intenta recuperar el texto tal cual como si hicieras Ctrl+C,
                        # evitando problemas de CID fonts raras.
                        paginas.append(page.get_text("text"))
                        
                texto_extraido = "\n".join(paginas)
                
            elif ext == '.docx':
                if not Document:
                    raise ImportError("Librería python-docx no instalada.")
                doc = Document(file_obj)
                parrafos = [p.text for p in doc.paragraphs if p.text.strip() != ""]
                texto_extraido = "\n".join(parrafos)
                
            elif ext == '.pptx':
                if not Presentation:
                    raise ImportError("Librería python-pptx no instalada.")
                prs = Presentation(file_obj)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                texto_extraido = "\n".join(text_runs)
                
            elif ext in ['.csv', '.xlsx', '.xls']:
                # Extraemos todo el texto de todas las celdas como un gran corpus
                if ext == '.csv':
                    df = pd.read_csv(file_obj)
                else:
                    df = pd.read_excel(file_obj)
                
                # Convertimos todo a string y lo unimos
                text_series = df.astype(str).apply(lambda x: ' '.join(x), axis=1)
                texto_extraido = "\n".join(text_series.tolist())

            elif ext in ['.mp3', '.wav', '.m4a', '.mp4']:
                if not whisper:
                    raise ImportError("Librería Whisper no instalada o falla al cargar.")
                    
                import tempfile
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_audio:
                    tmp_audio.write(file_obj.getvalue())
                    tmp_audio_path = tmp_audio.name
                
                try:
                    # El modelo base requiere al menos 1 GB de VRAM/RAM.
                    model = whisper.load_model("base") 
                    result = model.transcribe(tmp_audio_path)
                    texto_extraido = result["text"]
                except Exception as e:
                    texto_extraido = f"[Error Transcripción Whisper. Se requiere FFMPEG instalado en el SO: {str(e)}]"
                finally:
                    os.remove(tmp_audio_path)
                
            else:
                raise ValueError(f"Formato no soportado: {ext}")
                
        except Exception as e:
            texto_extraido = f"[Error al extraer texto de {filename}: {str(e)}]"

        return texto_extraido
